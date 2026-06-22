#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Generate the group-meeting presentation deck.

Structure follows the whiteboard outline:
  1. Motivation
  2. Pipeline (overview + AgentClinic base system)
  3. Version Update Detection (categorical + embedding: full conv / diagnosis only)
  4. Evaluation (history borrowing vs non-borrowing, vs ground truth)

Figures/results are left as labeled placeholders to be filled in later.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ----------------------------------------------------------------------------
# Theme
# ----------------------------------------------------------------------------
NAVY   = RGBColor(0x12, 0x2A, 0x4A)   # deep navy   - titles / bars
BLUE   = RGBColor(0x2F, 0x6F, 0xB5)   # accent blue
TEAL   = RGBColor(0x1F, 0x9E, 0x8F)   # accent teal
GREY   = RGBColor(0x55, 0x5B, 0x66)   # body grey
LGREY  = RGBColor(0x8A, 0x90, 0x99)   # light grey
BG     = RGBColor(0xFF, 0xFF, 0xFF)
PH_BG  = RGBColor(0xF1, 0xF4, 0xF8)   # placeholder fill
PH_LN  = RGBColor(0xB6, 0xC2, 0xD2)   # placeholder border

FONT = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------------
def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element)
    s.shapes._spTree.insert(2, bg._element)
    return s


def _set(tf, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT,
         font=FONT, space_after=6):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.name = font; f.color.rgb = color
    return p


def textbox(s, left, top, width, height, text, size, color, **kw):
    tb = s.shapes.add_textbox(left, top, width, height)
    _set(tb.text_frame, text, size, color, **kw)
    return tb


def bullets(s, left, top, width, height, items, size=18,
            color=GREY, leading_color=NAVY):
    """items: list of (level, text, bold) ; level 0/1/2"""
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for lvl, txt, bold in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(8 if lvl == 0 else 4)
        p.space_before = Pt(4 if lvl == 0 else 0)
        bullet = "—  " if lvl == 0 else ("•  " if lvl == 1 else "·  ")
        r = p.add_run(); r.text = bullet + txt
        f = r.font
        f.name = FONT
        f.size = Pt(size if lvl == 0 else size - 2 - lvl)
        f.bold = bold or (lvl == 0)
        f.color.rgb = leading_color if lvl == 0 else color
    return tb


def header(s, title, kicker=None):
    # top accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background(); bar.shadow.inherit = False
    if kicker:
        textbox(s, Inches(0.6), Inches(0.14), Inches(11), Inches(0.3),
                kicker.upper(), 12, RGBColor(0x9C, 0xC2, 0xEC), bold=True)
    textbox(s, Inches(0.6), Inches(0.40), Inches(12.2), Inches(0.7),
            title, 28, BG, bold=True)
    # thin teal underline accent
    ac = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.15),
                            Inches(2.2), Inches(0.06))
    ac.fill.solid(); ac.fill.fore_color.rgb = TEAL
    ac.line.fill.background(); ac.shadow.inherit = False


def placeholder(s, left, top, width, height, label):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = PH_BG
    box.line.color.rgb = PH_LN; box.line.width = Pt(1.25)
    box.line.dash_style = 2  # dashed
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "▣  " + label
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = LGREY
    r.font.name = FONT
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = "placeholder — drop figure/table here"
    r2.font.size = Pt(10); r2.font.italic = True; r2.font.color.rgb = LGREY
    r2.font.name = FONT
    return box


def chip(s, left, top, width, text, fill, txt_color=BG, height=Inches(0.55)):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.fill.background(); c.shadow.inherit = False
    tf = c.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = txt_color; r.font.name = FONT
    return c


def arrow(s, left, top, width, height=Inches(0.4), color=BLUE):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background(); a.shadow.inherit = False
    return a


def pagenum(s, n):
    textbox(s, Inches(12.3), Inches(7.05), Inches(0.9), Inches(0.35),
            str(n), 11, LGREY, align=PP_ALIGN.RIGHT)
    textbox(s, Inches(0.6), Inches(7.05), Inches(8), Inches(0.35),
            "Version-Aware Clinical Trial  ·  Group Meeting 12/19", 10, LGREY)


# ============================================================================
# SLIDE 1 — Title
# ============================================================================
s = slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.15), SW, Inches(2.9))
band.fill.solid(); band.fill.fore_color.rgb = NAVY
band.line.fill.background(); band.shadow.inherit = False
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.15), SW, Inches(0.09))
acc.fill.solid(); acc.fill.fore_color.rgb = TEAL
acc.line.fill.background(); acc.shadow.inherit = False

textbox(s, Inches(0.9), Inches(0.9), Inches(11.5), Inches(0.5),
        "Group Meeting · Summer Project", 16, BLUE, bold=True)
textbox(s, Inches(0.9), Inches(2.45), Inches(11.6), Inches(1.4),
        "Version-Aware Evaluation of LLM Clinical Agents", 40, BG, bold=True)
textbox(s, Inches(0.9), Inches(3.75), Inches(11.6), Inches(0.8),
        "Detecting silent model updates and estimating performance with history borrowing",
        18, RGBColor(0xCF, 0xDD, 0xEE))
textbox(s, Inches(0.9), Inches(5.5), Inches(11.5), Inches(0.5),
        "Built on the AgentClinic simulation framework", 15, GREY, italic=True)
textbox(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5),
        "Room 402C  ·  Friday 12/19  ·  12–1 pm", 14, LGREY)

# ============================================================================
# SLIDE 2 — Agenda
# ============================================================================
s = slide(); header(s, "Agenda", kicker="What we'll cover")
items = [
    ("1", "Motivation", "Why version-aware evaluation matters for clinical LLM agents"),
    ("2", "Pipeline", "The AgentClinic base system and our trial layer"),
    ("3", "Version Update Detection", "Categorical and embedding-based signals"),
    ("4", "Evaluation", "History borrowing vs. non-borrowing against ground truth"),
]
y = 1.6
for num, t, d in items:
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y), Inches(0.7), Inches(0.7))
    c.fill.solid(); c.fill.fore_color.rgb = TEAL; c.line.fill.background(); c.shadow.inherit = False
    tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.bold = True; r.font.size = Pt(20); r.font.color.rgb = BG
    textbox(s, Inches(1.75), Inches(y - 0.02), Inches(5), Inches(0.5), t, 22, NAVY, bold=True)
    textbox(s, Inches(1.75), Inches(y + 0.42), Inches(10.5), Inches(0.4), d, 14, GREY)
    y += 1.25
pagenum(s, 2)

# ============================================================================
# SLIDE 3 — Motivation
# ============================================================================
s = slide(); header(s, "Motivation", kicker="Section 1")
bullets(s, Inches(0.7), Inches(1.5), Inches(7.2), Inches(5.4), [
    (0, "LLMs behind clinical agents change silently", True),
    (1, "Providers update models and prompts continuously (GPT, DeepSeek, ...)", False),
    (1, "The same agent can behave differently week to week", False),
    (0, "A full re-trial after every update is infeasible", True),
    (1, "Clinical-style evaluation is slow and expensive per case", False),
    (1, "We cannot run a fresh RCT every time a version ships", False),
    (0, "We need two capabilities", True),
    (2, "Detect when an update is meaningful enough to matter", False),
    (2, "Estimate the new version's performance efficiently", False),
    (0, "Idea: treat the model version as a time-varying intervention", True),
    (1, "Borrow historical evaluation data instead of starting from zero", False),
])
placeholder(s, Inches(8.2), Inches(1.6), Inches(4.5), Inches(4.6),
            "Motivating example: accuracy drift across versions")
pagenum(s, 3)

# ============================================================================
# SLIDE 4 — Pipeline overview (big picture)
# ============================================================================
s = slide(); header(s, "Pipeline Overview", kicker="Section 2 · The big picture")
# flow chips
top = Inches(1.7)
chip(s, Inches(0.6), top, Inches(2.5), "Case stream\n(MedQA / NEJM)", BLUE)
arrow(s, Inches(3.2), Inches(1.78), Inches(0.7))
chip(s, Inches(4.0), top, Inches(2.6), "AgentClinic\nsimulation", NAVY)
arrow(s, Inches(6.75), Inches(1.78), Inches(0.7))
chip(s, Inches(7.55), top, Inches(2.5), "Dialogues +\ndiagnoses", TEAL)
arrow(s, Inches(10.2), Inches(1.78), Inches(0.7))
chip(s, Inches(11.0), top, Inches(1.9), "Trial log\n(JSONL)", GREY)

textbox(s, Inches(0.6), Inches(2.75), Inches(12), Inches(0.4),
        "From the logged data we run two downstream analyses:", 15, NAVY, bold=True)

placeholder(s, Inches(0.6), Inches(3.3), Inches(5.9), Inches(2.4),
            "Big-picture system diagram (overview)")
bullets(s, Inches(6.9), Inches(3.25), Inches(6.0), Inches(2.7), [
    (0, "Version Update Detection", True),
    (1, "Did the model/prompt change meaningfully?", False),
    (0, "Performance Estimation", True),
    (1, "How accurate is the new version, using less data?", False),
    (0, "Base system: AgentClinic", True),
    (1, "Models tested: DeepSeek, GPT, Qwen", False),
])
pagenum(s, 4)

# ============================================================================
# SLIDE 5 — AgentClinic base system
# ============================================================================
s = slide(); header(s, "Base System: AgentClinic", kicker="Section 2 · Detail")
bullets(s, Inches(0.7), Inches(1.45), Inches(6.4), Inches(5.5), [
    (0, "Multi-agent clinical simulation", True),
    (1, "Doctor agent: interviews, orders tests, diagnoses", False),
    (1, "Patient agent: answers from a hidden case profile", False),
    (1, "Measurement agent: returns test/exam results", False),
    (1, "Moderator: judges diagnosis vs. ground truth", False),
    (0, "Two-doctor consensus workflow", True),
    (1, "Reviewer agent re-diagnoses from the transcript", False),
    (1, "Agree → scored against ground truth", False),
    (1, "Disagree → excluded; discussion logged", False),
    (0, "Our trial layer on top", True),
    (1, "Version epochs, 1:1 control/treatment randomization", False),
    (1, "Append-only JSONL logging per case", False),
])
# agent loop diagram (simple)
placeholder(s, Inches(7.4), Inches(1.5), Inches(5.3), Inches(2.6),
            "Doctor ⇄ Patient / Measurement loop")
placeholder(s, Inches(7.4), Inches(4.3), Inches(5.3), Inches(1.9),
            "Two-doctor consensus → accuracy")
pagenum(s, 5)

# ============================================================================
# SLIDE 6 — Data & models
# ============================================================================
s = slide(); header(s, "Data & Models", kicker="Section 2 · Setup")
bullets(s, Inches(0.7), Inches(1.45), Inches(6.2), Inches(5.4), [
    (0, "Datasets", True),
    (1, "MedQA / NEJM clinical scenarios", False),
    (1, "Ground-truth diagnosis per case", False),
    (0, "Models under test", True),
    (1, "DeepSeek: flash, pro", False),
    (1, "GPT: gpt-5, gpt-5-mini, gpt-5.5", False),
    (1, "Qwen: plus / turbo (embedding comparison)", False),
    (0, "Replication", True),
    (1, "Multiple independent runs per model per case", False),
    (1, "Enables per-case distribution + variance analysis", False),
])
placeholder(s, Inches(7.3), Inches(1.5), Inches(5.4), Inches(4.7),
            "Data summary: #cases × #models × #replicates")
pagenum(s, 6)

# ============================================================================
# SLIDE 7 — Version Update Detection overview
# ============================================================================
s = slide(); header(s, "Version Update Detection", kicker="Section 3 · Overview")
textbox(s, Inches(0.7), Inches(1.35), Inches(12), Inches(0.5),
        "Given a candidate (LLM, prompt), decide whether behavior changed meaningfully.",
        16, NAVY, bold=True)
# two columns
chip(s, Inches(0.9), Inches(2.1), Inches(5.2), "Approach A — Categorical", BLUE, height=Inches(0.6))
bullets(s, Inches(0.9), Inches(2.95), Inches(5.3), Inches(4), [
    (1, "Map diagnoses to ICD-10 categories", False),
    (1, "Build per-model diagnosis distribution", False),
    (1, "Compare with JS / symmetric-KL divergence", False),
    (1, "Interpretable, discrete shift signal", False),
])
chip(s, Inches(6.9), Inches(2.1), Inches(5.6), "Approach B — Embedding", TEAL, height=Inches(0.6))
bullets(s, Inches(6.9), Inches(2.95), Inches(5.6), Inches(4), [
    (1, "Embed text, compare cosine similarity", False),
    (2, "Full conversation  vs.  diagnosis-only", False),
    (1, "Case-level + model-level similarity", False),
    (1, "Surfaces low-similarity (changed) cases", False),
])
pagenum(s, 7)

# ============================================================================
# SLIDE 8 — Categorical detection
# ============================================================================
s = slide(); header(s, "Detection A — Categorical (ICD-10)", kicker="Section 3")
bullets(s, Inches(0.7), Inches(1.45), Inches(6.0), Inches(5.2), [
    (0, "Pipeline", True),
    (1, "Normalize each diagnosis to an ICD-10 code", False),
    (1, "Aggregate into a category distribution per model", False),
    (1, "Distance = Jensen–Shannon / symmetric KL", False),
    (0, "Reading the matrix", True),
    (1, "Within-model replicates → baseline noise floor", False),
    (1, "Across models → real categorical shift", False),
    (0, "Use case", True),
    (1, "Cheap first-pass flag for distributional drift", False),
])
placeholder(s, Inches(7.0), Inches(1.5), Inches(5.7), Inches(4.7),
            "ICD-10 distribution divergence heatmap")
pagenum(s, 8)

# ============================================================================
# SLIDE 9 — Embedding detection
# ============================================================================
s = slide(); header(s, "Detection B — Embedding Similarity", kicker="Section 3")
bullets(s, Inches(0.7), Inches(1.4), Inches(5.6), Inches(5.4), [
    (0, "Two granularities", True),
    (1, "Full conversation: whole doctor–patient transcript", False),
    (1, "Diagnosis-only: final diagnosis text", False),
    (0, "Procedure", True),
    (1, "Embed → average within case/group → normalize", False),
    (1, "Cosine similarity across model pairs, per case", False),
    (1, "Average across cases → model-level similarity", False),
    (0, "Outputs", True),
    (1, "Similarity matrix + low-similarity case lists", False),
])
placeholder(s, Inches(6.6), Inches(1.5), Inches(6.1), Inches(2.55),
            "Model-level similarity matrix (full conv.)")
placeholder(s, Inches(6.6), Inches(4.2), Inches(6.1), Inches(2.0),
            "Full conv. vs. diagnosis-only comparison")
pagenum(s, 9)

# ============================================================================
# SLIDE 10 — Detection takeaways
# ============================================================================
s = slide(); header(s, "Detection — Preliminary Findings", kicker="Section 3 · Results")
bullets(s, Inches(0.7), Inches(1.45), Inches(6.2), Inches(5.2), [
    (0, "Similarity tracks model closeness", True),
    (1, "flash vs. pro most similar (~0.98 full-dialogue)", False),
    (1, "Cross-family pairs (e.g. qwen vs. gpt) lowest", False),
    (0, "Full conversation > diagnosis-only signal", True),
    (1, "Whole transcript captures behavioral change earlier", False),
    (0, "Open question", True),
    (1, "Where to set the 'meaningful update' threshold?", False),
])
placeholder(s, Inches(7.2), Inches(1.5), Inches(5.5), Inches(4.7),
            "Key result figure (replace numbers above)")
pagenum(s, 10)

# ============================================================================
# SLIDE 11 — Evaluation setup
# ============================================================================
s = slide(); header(s, "Evaluation", kicker="Section 4 · Setup")
bullets(s, Inches(0.7), Inches(1.45), Inches(6.2), Inches(5.2), [
    (0, "Question", True),
    (1, "Can we estimate a version's accuracy with less data", False),
    (1, "by borrowing from historical / similar models?", False),
    (0, "Models evaluated", True),
    (1, "gpt-5, gpt-5-mini, gpt-5.5", False),
    (1, "DeepSeek flash, DeepSeek pro", False),
    (0, "Comparison", True),
    (1, "History borrowing  vs.  non-borrowing", False),
    (1, "Both measured against full ground-truth accuracy", False),
])
placeholder(s, Inches(7.2), Inches(1.5), Inches(5.5), Inches(4.7),
            "Evaluation design diagram")
pagenum(s, 11)

# ============================================================================
# SLIDE 12 — History borrowing method
# ============================================================================
s = slide(); header(s, "Method — Similarity-Weighted History Borrowing", kicker="Section 4")
# formula box
fb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.45),
                        Inches(12), Inches(1.0))
fb.fill.solid(); fb.fill.fore_color.rgb = RGBColor(0xED, 0xF3, 0xFA)
fb.line.color.rgb = BLUE; fb.line.width = Pt(1); fb.shadow.inherit = False
tf = fb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "θ_borrowed(j) = α · θ_small(j) + (1 − α) · Σ_i  w_ij · θ_small(i)        w_ij ∝ exp(−λ · d(i,j))"
r.font.size = Pt(17); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = NAVY
bullets(s, Inches(0.7), Inches(2.8), Inches(6.1), Inches(4), [
    (0, "Ingredients", True),
    (1, "θ_small: accuracy from a small case budget", False),
    (1, "d(i,j) = 1 − similarity (from Section 3)", False),
    (1, "α balances self vs. borrowed evidence", False),
    (1, "λ controls how fast distant models are down-weighted", False),
    (0, "Training", True),
    (1, "Grid-search one shared (α, λ) across orders & models", False),
    (1, "Minimize MAE vs. ground-truth accuracy", False),
])
placeholder(s, Inches(7.0), Inches(2.85), Inches(5.7), Inches(3.6),
            "Borrowing schematic / (α, λ) landscape")
pagenum(s, 12)

# ============================================================================
# SLIDE 13 — Results: borrowing vs non-borrowing
# ============================================================================
s = slide(); header(s, "Results — Borrowing vs. Non-Borrowing", kicker="Section 4 · Results")
placeholder(s, Inches(0.6), Inches(1.5), Inches(7.3), Inches(4.9),
            "MAE vs. case budget: borrowing vs. non-borrowing")
bullets(s, Inches(8.2), Inches(1.6), Inches(4.6), Inches(4.8), [
    (0, "What to look for", True),
    (1, "Borrowing lowers error at small budgets", False),
    (1, "Gap shrinks as budget grows", False),
    (0, "Per-model view", True),
    (1, "Largest gains for models with close peers", False),
    (0, "Takeaway", True),
    (1, "Similar history → fewer cases for same accuracy", False),
])
pagenum(s, 13)

# ============================================================================
# SLIDE 14 — Ground truth comparison
# ============================================================================
s = slide(); header(s, "Validation Against Ground Truth", kicker="Section 4 · Results")
placeholder(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(4.9),
            "Estimated vs. ground-truth accuracy (per model)")
placeholder(s, Inches(6.8), Inches(1.5), Inches(5.9), Inches(4.9),
            "Error table: borrow vs. non-borrow vs. truth")
pagenum(s, 14)

# ============================================================================
# SLIDE 15 — Summary & next steps
# ============================================================================
s = slide(); header(s, "Summary & Next Steps", kicker="Wrap-up")
bullets(s, Inches(0.7), Inches(1.45), Inches(6.1), Inches(5.4), [
    (0, "Summary", True),
    (1, "AgentClinic-based version-aware trial pipeline", False),
    (1, "Two detection signals: categorical + embedding", False),
    (1, "History borrowing improves small-budget estimates", False),
    (0, "Next steps", True),
    (1, "Principled update threshold from detection signals", False),
    (1, "Time-discounted / drift-triggered borrowing", False),
    (1, "Broader outcomes: safety, compliance, tool use", False),
])
placeholder(s, Inches(7.1), Inches(1.5), Inches(5.6), Inches(4.7),
            "Roadmap timeline")
pagenum(s, 15)

# ============================================================================
# SLIDE 16 — Discussion / Q&A
# ============================================================================
s = slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
band.fill.solid(); band.fill.fore_color.rgb = NAVY
band.line.fill.background(); band.shadow.inherit = False
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.6), SW, Inches(0.09))
acc.fill.solid(); acc.fill.fore_color.rgb = TEAL
acc.line.fill.background(); acc.shadow.inherit = False
textbox(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.0),
        "Discussion & Questions", 40, BG, bold=True)
textbox(s, Inches(0.9), Inches(3.9), Inches(11.5), Inches(0.8),
        "Feedback welcome — especially on detection thresholds and borrowing assumptions.",
        17, RGBColor(0xCF, 0xDD, 0xEE))
textbox(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.5),
        "Thank you!", 16, LGREY)

out = "Version_Aware_Clinical_Trial.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
